#%%
#!/usr/bin/env python3
"""
MIT License

Copyright (c) 2023 Okyaz Eminaga

Permission is hereby granted, free of charge, to any person obtaining a copy
of this software and associated documentation files (the "Software"), to deal
in the Software without restriction, including without limitation the rights
to use, copy, modify, merge, publish, distribute, sublicense, and/or sell
copies of the Software, and to permit persons to whom the Software is
furnished to do so, subject to the following conditions:

The above copyright notice and this permission notice shall be included in all
copies or substantial portions of the Software.

THE SOFTWARE IS PROVIDED "AS IS", WITHOUT WARRANTY OF ANY KIND, EXPRESS OR
IMPLIED, INCLUDING BUT NOT LIMITED TO THE WARRANTIES OF MERCHANTABILITY,
FITNESS FOR A PARTICULAR PURPOSE AND NONINFRINGEMENT. IN NO EVENT SHALL THE
AUTHORS OR COPYRIGHT HOLDERS BE LIABLE FOR ANY CLAIM, DAMAGES OR OTHER
LIABILITY, WHETHER IN AN ACTION OF CONTRACT, TORT OR OTHERWISE, ARISING FROM,
OUT OF OR IN CONNECTION WITH THE SOFTWARE OR THE USE OR OTHER DEALINGS IN THE
SOFTWARE.
"""
from pptx import Presentation
from pptx.util import Inches
import pandas as pd
from tqdm import tqdm
import argparse
#%%
class GeneratePowerPointPresentation(object):
    def __init__(self, csv_file, destination) -> None:
        self.csv_file = csv_file
        self.destination =destination
        self.data = pd.read_csv(self.csv_file)

    def CreateAndSave(self, filename,rows, destination):
        prs = Presentation()
        ht = {"WLC":[], "BLC": []}
        others = []
        rows=rows.sort_values(by="IDENTIFICATION")
        idx = list(rows.IDENTIFICATION.unique())
        for i in idx:
            ht = {"WLC":[], "BLC": [], "H":[]}
            others = []
            selected_itms =rows[rows.IDENTIFICATION==i]
            for _, itm in selected_itms.iterrows():
                key = str(itm["ImageModality"])
                if "H" in key:
                    ht["H"].append(itm)
                elif key == "WLC":
                    ht["WLC"].append(itm)
                elif key == "BLC":
                    ht["BLC"].append(itm)
                else:
                    others.append(itm)

            for key in ht:
                rows__ = ht[key]
                for r_ in rows__:
                    try:
                        prs=self.CreateSlide(prs, r_)
                    except:
                        print("ERR: Can't generate a slide for the row: ",r_)

            for _r_ in others:
                try:
                    prs=self.CreateSlide(prs, _r_)
                except:
                    print("ERR: Can't generate a slide for the row: ",_r_)
        prs.save(f'{destination}/{filename}.pptx')
        
    def CreateSlide(self, prs, row):
        CML = row["ID"]
        DATE = row["DATE"]
        IND = "TURBT" if row["IND"]=="O" else "Office Cystoscopy"
        TYPE_OF_IMAGE = row["ImageModality"]
        IDENTIFICATION=row["IDENTIFICATION"]
        LOCATION= row["LOCATION"]
        PATHOLOGY= row["PATHOLOGY"]
        STAGE=row["STAGE"]

        criteria =[TYPE_OF_IMAGE, IDENTIFICATION,LOCATION, PATHOLOGY,STAGE]
        for i, cri in enumerate(criteria):
            if str(cri)== "nan":
                criteria[i]= ""

        TYPE_OF_IMAGE=criteria[0]
        IDENTIFICATION=criteria[1]
        LOCATION=criteria[2]
        PATHOLOGY=criteria[3]
        STAGE=criteria[4]
        if LOCATION == "":
            DESCRIPTION = f" {IDENTIFICATION}"
        else:
            DESCRIPTION = f"Lesion id: {IDENTIFICATION} | Image modality: {TYPE_OF_IMAGE} | Location: {LOCATION} | Pathology: {STAGE} {PATHOLOGY}"
        IMG_PATH = row["FullPath"]
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        left = Inches(0.5)
        top = Inches(0.14)
        width = Inches(8.5)
        height = Inches(0.4)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = f"{CML} - {DATE}" # The information from the filename should be presented here. We may link also REDCap. 
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(8.5)
        height = Inches(0.8)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = f"{DESCRIPTION}"
        img_path = f"{IMG_PATH}" #Filename of the screenshots
        left = Inches(0.5)
        top = Inches(1.33)
        height = Inches(9)
        pic = slide.shapes.add_picture(img_path, left, top, height)
        return prs

    def Run(self):
        list_of_unique_cases =list(self.data.CASE_ID.unique())

        for case_id in tqdm(list_of_unique_cases):
            selected_data = self.data[self.data.CASE_ID==case_id]
            self.CreateAndSave(case_id, selected_data, self.destination)
# %%
parser = argparse.ArgumentParser(description='User Interface for image labelling')
parser.add_argument("--csv", help="provide the csv file with a column FullPath.")
parser.add_argument("--dst", help="provide the folder to which the power point presentations will be stored.")

args = parser.parse_args()

if __name__ == "__main__":
    print(args)
    cmd=GeneratePowerPointPresentation(args.csv,args.dst)
    cmd.Run()